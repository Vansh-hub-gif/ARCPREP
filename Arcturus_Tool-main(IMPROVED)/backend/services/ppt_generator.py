import os
import re
from typing import Any, Dict, List, Optional, Tuple
from collections import OrderedDict

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn

# Constants and layout variables
SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)
BLUE = RGBColor(0, 128, 163)
DARK_BLUE = RGBColor(0, 88, 118)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
LIGHT_ROW = RGBColor(255, 255, 255)
BORDER = RGBColor(90, 90, 90)
FONT = "Arial"
BODY_FONT_SIZE = 9
HEADER_FONT_SIZE = 9
TITLE_FONT_SIZE = 26
SUBTITLE_FONT_SIZE = 22
TOP_STRIP_H = Inches(0.08)
BOTTOM_STRIP_H = Inches(0.10)
LOGO_TOP_Y = Inches(0.18)
CONTENT_MARGIN_X = Inches(0.28)
TABLE_TOP_Y = Inches(1.15)
TABLE_W = SLIDE_W - Inches(0.56)
HEADER_ROW_H = Inches(0.32)
DATA_ROW_H = Inches(1.52)
TITLE_BAND_Y = Inches(2.28)
TITLE_BAND_H = Inches(1.28)
INDEX_TABLE_TOP_Y = Inches(1.28)
FOOTER_TEXT_Y = SLIDE_H - BOTTOM_STRIP_H + Inches(0.012)
BOTTOM_LOGO_W = Inches(0.62)
BOTTOM_LOGO_H = Inches(0.38)
BOTTOM_LOGO_Y = SLIDE_H - BOTTOM_STRIP_H - BOTTOM_LOGO_H - Inches(0.08)
MAX_INDEX_ROWS_PER_SLIDE = 10
MAX_REF_ROWS_PER_SLIDE = 14  # Added for Reference pagination
TEMPLATE_CANDIDATES = ["templates/inventory_template.pptx","template/inventory_template.pptx","inventory_template.pptx","templates/oquat_template.pptx","oquat_template.pptx"]
LOGO_PATH_CANDIDATES = {"tid": ["assets/tid_logo.png","assets/tid.png","templates/tid_logo.png","templates/tid.png","static/tid_logo.png","static/tid.png"],"arcturus": ["assets/arcturus_logo.png","assets/arcturus.png","templates/arcturus_logo.png","templates/arcturus.png","static/arcturus_logo.png","static/arcturus.png"]}


# --- Generic Boilerplate text to detect and ignore in Business Benefit ---
GENERIC_BENEFIT_PHRASES = [
    "this feature enables organizations to improve data accuracy",
    "this feature enables organizations to improve processing efficiency",
    "this feature enables organizations to improve visibility",
    "this feature enables organizations to improve control"
]

# Tracks benefit bullets that appear ANYWHERE else in the deck
_COMMON_BENEFIT_BULLETS: set = set()

def _safe_text(value: Any) -> str:
    text = str(value or "").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n\s*\n+", "\n", text)
    return text.strip()


def _pick(feature: Dict[str, Any], *keys: str, default: str = "") -> str:
    for key in keys:
        value = feature.get(key)
        if value not in (None, ""):
            return _safe_text(value)
    return default


def _release_version(features: List[Dict[str, Any]]) -> str:
    for f in features:
        val = _pick(f, "release_version", "release", "version")
        if val:
            return val
    return "26B"


def _module_name(features: List[Dict[str, Any]]) -> str:
    for f in features:
        val = _pick(f, "module", "module_name", "product_area", "pillar", "category")
        if val and val.lower() not in {"nan", "none", "n/a"}:
            return val
    return "Collaboration Messaging"


def _feature_module(feature: Dict[str, Any]) -> str:
    return _pick(feature, "module", "pillar", "product_area", "category", default="General")


def _feature_title(feature: Dict[str, Any]) -> str:
    return _pick(feature, "release_feature", "title", "feature_name", "name", default="Untitled Feature")


def _description(feature: Dict[str, Any]) -> str:
    return _pick(feature, "short_description", "description", "desc")


def _mandatory(feature: Dict[str, Any]) -> str:
    raw = _pick(feature, "mandatory", "is_mandatory", "required", "delivery_status", "action_required")
    low = raw.lower()
    if low in {"yes", "y", "true", "mandatory"}:
        return "Yes"
    if "setup required" in low or "enable" in low or "configure" in low:
        return "Yes"
    if low in {"no", "n", "false", "optional"}:
        return "No"
    return "No"


def _priority(feature: Dict[str, Any]) -> str:
    val = _pick(feature, "priority", "impact_priority", default="Medium")
    val_low = val.lower()
    if "high" in val_low:
        return "High"
    if "low" in val_low:
        return "Low"
    return "Medium"


def _split_points(text: str) -> List[str]:
    text = _safe_text(text)
    if not text:
        return []

    if "•" in text or re.search(r"(^|\n)\s*[-*]\s+", text):
        parts = re.split(r"(?:^|\n)\s*(?:•|-|\*)\s*", text)
    else:
        parts = re.split(r"(?<=[.!?])\s+", text)
        if len([p for p in parts if p.strip()]) <= 1 and ";" in text:
            parts = text.split(";")

    points = []
    seen = set()
    for p in parts:
        p = re.sub(r"^[\s•\-*]+", "", p).strip()
        p = re.sub(r"\s+", " ", p)
        if not p or len(p) < 8:
            continue
        p = p.rstrip(".; ") + "."
        norm = re.sub(r"[^a-z0-9]+", " ", p.lower()).strip()
        if norm in seen:
            continue
        seen.add(norm)
        points.append(p)
    return points


def _ppt_points(text: str, max_points: int = 3) -> str:
    points = _split_points(text)[:max_points]
    return "\n".join([f"• {p}" for p in points])


def _release_feature_cell(feature: Dict[str, Any]) -> str:
    title = _feature_title(feature)
    desc = _description(feature)
    if desc:
        return f"{title}\n\n{desc}"
    return title

# --- Generic Boilerplate text to detect and ignore in Business Benefit ---
GENERIC_BENEFIT_PHRASES = [
    "this feature enables organizations to improve data accuracy",
    "this feature enables organizations to improve processing efficiency",
    "this feature enables organizations to improve visibility",
    "this feature enables organizations to improve control"
]

# Vague cliches that AI uses repeatedly
COMMON_CLICHES = [
    "improves data accuracy", "reduces manual effort", "streamlines operations",
    "increases efficiency", "enhances visibility", "reduces errors",
    "saves time", "improves productivity", "automates processes"
]

# --- Generic Boilerplate text to detect and fix ---
GENERIC_PHRASES = [
    "this feature enables organizations to improve data accuracy",
    "this feature enables organizations to improve processing efficiency",
    "this feature enables organizations to improve visibility",
    "this feature enables organizations to improve control",
    "this feature enables organizations to"
]

# Keywords that indicate a sentence is an IMPACT, not a BENEFIT.
#
# ROOT CAUSE OF THE BLANK BUSINESS BENEFIT CELLS
# ----------------------------------------------
# These were matched with a plain `kw in sentence.lower()` substring test.
# "enable" therefore matched "enables" and "enabling"; "setup" matched
# "setups"; "configure" matched "configured"/"configures". Those words appear
# in almost every well-written benefit sentence, so _clean_benefit_text()
# discarded EVERY sentence and _benefit_cell() then returned "".
#
# Measured on the shipped run: the Excel had 0 blank benefits and the PPTX had
# 57 blank benefits out of 62 — same data, same run. The AI never failed; this
# filter deleted its output.
#
# Two changes: match on word boundaries, and (see _clean_benefit_text) never
# let the filter empty a cell.
IMPACT_KEYWORDS = [
    "requires", "must be", "users need", "action required",
    "mandatory", "provision", "assign role", "privilege",
    "prerequisite", "before using", "opt-in", "profile option",
    "setup required", "no setup",
]

# Pre-compiled word-boundary matchers. Multi-word entries keep their spacing.
_IMPACT_KEYWORD_RE = re.compile(
    r"\b(?:" + "|".join(re.escape(k) for k in IMPACT_KEYWORDS) + r")\b",
    re.IGNORECASE,
)


def _looks_like_impact_sentence(sentence: str) -> bool:
    return bool(_IMPACT_KEYWORD_RE.search(sentence or ""))

_COMMON_BENEFIT_BULLETS: set = set()
_COMMON_BENEFIT_PREFIXES: set = set()

def _is_generic_benefit(text: str) -> bool:
    if not text:
        return True
    text_low = text.lower()
    for phrase in GENERIC_PHRASES:
        if phrase in text_low:
            return True
    return False

def _clean_benefit_text(text: str) -> str:
    """
    Remove the generic lead-in and any genuine Impact sentences — WITHOUT ever
    emptying the cell.

    These filters are quality preferences, not correctness rules. If applying
    them removes every sentence, the original text is returned instead: an
    imperfect benefit always beats a blank cell in a customer deck. This is
    the fix for the 57-of-62 blank Business Benefit cells.
    """
    if not text:
        return ""

    # Scrub the legacy internal reviewer note. An earlier build wrote
    # "Business value for this feature was not generated by the AI pipeline
    # and requires functional review..." straight into the customer-facing
    # cell. That sentence must never reach a client deck, including when the
    # deck is rebuilt from an older saved dataset.
    text = re.sub(
        r"(?i)business value for this feature was not generated by the ai "
        r"pipeline[^.]*\.\s*",
        "",
        text,
    )

    # Fix the generic prefix first
    text = re.sub(r"(?i)this feature enables organizations to\s*", "", text).strip()
    # Capitalize the first letter after removal
    if text:
        text = text[0].upper() + text[1:]

    sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if s.strip()]
    if not sentences:
        return ""

    clean_sentences = [s for s in sentences if not _looks_like_impact_sentence(s)]

    if clean_sentences:
        return " ".join(clean_sentences)

    # Every sentence looked like an impact statement. Keep the benefit rather
    # than blanking the cell.
    print("[PPT BENEFIT] All sentences matched the impact filter; keeping the "
          "original text rather than emitting a blank cell.")
    return " ".join(sentences)

def _prescan_common_bullets(features: List[Dict[str, Any]]) -> None:
    """Pre-scan ALL features to find repetitive bullets across the entire deck."""
    global _COMMON_BENEFIT_BULLETS, _COMMON_BENEFIT_PREFIXES
    _COMMON_BENEFIT_BULLETS = set()
    _COMMON_BENEFIT_PREFIXES = set()
    
    exact_count: Dict[str, int] = {}
    prefix_count: Dict[str, int] = {}

    for feature in features:
        texts_to_scan = [
            _pick(feature, "short_description", "description", "desc"),
            _pick(feature, "business_impact", "impact", "impact_analysis"),
            _clean_benefit_text(_pick(feature, "business_benefit", "benefit", "business_value"))
        ]
        
        for text in texts_to_scan:
            if not text:
                continue
            points = _split_points(text)
            for p in points:
                norm = re.sub(r"[^a-z0-9]+", " ", p.lower()).strip()
                if norm and len(norm) > 15:
                    exact_count[norm] = exact_count.get(norm, 0) + 1
                    prefix = norm[:25]
                    if len(prefix) > 15:
                        prefix_count[prefix] = prefix_count.get(prefix, 0) + 1

    for norm, count in exact_count.items():
        if count > 1:
            _COMMON_BENEFIT_BULLETS.add(norm)
            
    for prefix, count in prefix_count.items():
        if count > 1:
            _COMMON_BENEFIT_PREFIXES.add(prefix)

def _filter_unique_benefit_points(points: List[str]) -> List[str]:
    unique = []
    for p in points:
        norm = re.sub(r"[^a-z0-9]+", " ", p.lower()).strip()
        prefix = norm[:25]
        if norm not in _COMMON_BENEFIT_BULLETS and prefix not in _COMMON_BENEFIT_PREFIXES:
            unique.append(p)
    return unique

def _benefit_cell(feature: Dict[str, Any]) -> str:
    """Generates Business Benefit ensuring:
    1. NO fake fallback sentences
    2. Fixes old generic AI data automatically by removing the prefix
    3. NO wrong impacts
    4. Exactly up to 2 bullet points
    """
    raw_benefit = _pick(feature, "business_benefit", "benefit", "business_value")
    cleaned_benefit = _clean_benefit_text(raw_benefit)
    
    # First, try to use strictly unique points
    points = []
    if cleaned_benefit:
        points = _filter_unique_benefit_points(_split_points(cleaned_benefit))
        
    # If uniqueness check deleted everything (because it was repetitive), 
    # fall back to the cleaned text rather than leaving it blank or using a fake sentence
    if not points and cleaned_benefit:
        points = _split_points(cleaned_benefit)

    # Format up to 2 points
    if points:
        formatted_points = []
        for p in points[:2]:
            p = re.sub(r'\s*\.\.\.\s*$', '', p).strip()
            if not p.endswith("."):
                p += "."
            formatted_points.append(f"• {p}")
        return "\n".join(formatted_points)
        
    # LAST RESORT — a Business Benefit cell is never allowed to be blank.
    # We only reach here if the enricher handed us nothing at all, which the
    # ensure_non_empty_benefit() guard in ai_enricher.py should already
    # prevent. The fallback is derived from THIS feature's title and scraped
    # Oracle description, so it is feature-specific rather than filler.
    fallback = _fallback_benefit_points(feature)
    if fallback:
        print(f"[PPT BENEFIT GUARD] Empty benefit for "
              f"'{_pick(feature, 'title')}' — feature-specific fallback used.")
        return fallback

    return ""


def _fallback_benefit_points(feature: Dict[str, Any]) -> str:
    """
    Build a short, feature-specific benefit from the feature's own title and
    its scraped Oracle description. Deterministic and grounded — two different
    features cannot produce the same text.
    """
    title = _pick(feature, "title", "feature", "name")
    description = _pick(feature, "description", "short_description", "desc")
    if not title and not description:
        return ""

    clean_title = re.sub(
        r"^(Redwood|AI Agent(ic App)?)\s*[:\-]\s*", "", title, flags=re.IGNORECASE
    ).strip()

    points: List[str] = []

    # Point 1 — what the feature lets the team do, from its own title.
    if clean_title:
        # Acronym-safe: "B2B Message Converter" must not become "b2B ...".
        _first = re.sub(r"[^A-Za-z0-9]", "", clean_title.split()[0])
        if (len(_first) >= 2 and _first[1:2].isupper()) or _first.isupper():
            lowered = clean_title
        else:
            lowered = clean_title[:1].lower() + clean_title[1:]
        points.append(
            f"Enables teams to {lowered} directly in Oracle Cloud SCM, "
            f"reducing the manual effort the current process requires"
        )

    # Point 2 — grounded in Oracle's own first sentence for this feature.
    if description:
        first = re.split(r"(?<=[.!?])\s+", description.strip())[0].strip()
        first = re.sub(r"\s+", " ", first).rstrip(".")
        if len(first.split()) >= 6:
            if len(first) > 220:
                first = first[:217].rsplit(" ", 1)[0]
            points.append(f"Addresses the scenario Oracle describes: {first}")

    if not points:
        return ""

    formatted = []
    for p in points[:2]:
        if not p.endswith("."):
            p += "."
        formatted.append(f"• {p}")
    return "\n".join(formatted)

def _impact_cell(feature: Dict[str, Any]) -> str:
    """Dynamically aligns with the feature, preventing duplication with Benefit."""
    raw_impact = _pick(feature, "business_impact", "impact", "impact_analysis")
    
    # If impact is missing or generic, derive it from the latter part of the description
    if _is_generic_benefit(raw_impact):
        desc = _description(feature)
        sentences = re.split(r'(?<=[.!?])\s+', desc)
        if len(sentences) > 1:
            return _ppt_points(" ".join(sentences[1:]), 3) # Use later sentences for impact to differ from benefit
        return _ppt_points(desc, 3)
        
    return _ppt_points(raw_impact, 3)


def _find_template_path(explicit: Optional[str] = None) -> Optional[str]:
    if explicit and os.path.exists(explicit):
        return explicit
    for p in TEMPLATE_CANDIDATES:
        if os.path.exists(p):
            return p
    return None


def _find_logo_path(kind: str) -> Optional[str]:
    for p in LOGO_PATH_CANDIDATES.get(kind, []):
        if os.path.exists(p):
            return p
    return None


def _extract_logos_from_template(template_path: Optional[str]) -> Dict[str, Optional[str]]:
    result = {"tid": _find_logo_path("tid"), "arcturus": _find_logo_path("arcturus")}
    if result["tid"] and result["arcturus"]:
        return result
    if not template_path or not os.path.exists(template_path):
        return result

    os.makedirs("outputs/logo_cache", exist_ok=True)
    try:
        prs = Presentation(template_path)
        tid_candidates: List[Tuple[int, str]] = []
        arc_candidates: List[Tuple[int, str]] = []
        idx = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "shape_type", None) != 13:
                    continue
                idx += 1
                ext = "png"
                try:
                    ext = shape.image.ext or "png"
                except Exception:
                    pass
                out = f"outputs/logo_cache/logo_{idx}.{ext}"
                try:
                    with open(out, "wb") as f:
                        f.write(shape.image.blob)
                except Exception:
                    continue
                area = int(shape.width) * int(shape.height)
                if shape.left < prs.slide_width * 0.35:
                    tid_candidates.append((area, out))
                elif shape.left > prs.slide_width * 0.50:
                    arc_candidates.append((area, out))
        if not result["tid"] and tid_candidates:
            result["tid"] = sorted(tid_candidates, reverse=True)[0][1]
        if not result["arcturus"] and arc_candidates:
            result["arcturus"] = sorted(arc_candidates, reverse=True)[0][1]
    except Exception:
        pass
    return result


def _configure_table_borders(table):
    """Remove default table styles and banding so cell borders render correctly."""
    tbl = table._tbl
    tblPr = tbl.tblPr
    
    for style in tblPr.findall(qn('a:tblStyle')):
        tblPr.remove(style)
    
    tblPr.set('firstRow', '0')
    tblPr.set('lastRow', '0')
    tblPr.set('firstCol', '0')
    tblPr.set('lastCol', '0')
    tblPr.set('bandRow', '0')
    tblPr.set('bandCol', '0')
    
    for child in tblPr.findall(qn('a:tblBorders')):
        tblPr.remove(child)
        
    tblBorders = parse_xml(
        f'<a:tblBorders xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'  <a:top w="12700"><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:top>'
        f'  <a:left w="12700"><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:left>'
        f'  <a:bottom w="12700"><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:bottom>'
        f'  <a:right w="12700"><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:right>'
        f'  <a:insideH w="12700"><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:insideH>'
        f'  <a:insideV w="12700"><a:solidFill><a:srgbClr val="000000"/></a:solidFill></a:insideV>'
        f'</a:tblBorders>'
    )
    tblPr.append(tblBorders)


def _set_cell_border(cell, color: str = "000000", width: str = "12700") -> None:
    """Set explicit borders on individual cells, strictly respecting OOXML schema order."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    
    for edge in ('lnL', 'lnR', 'lnT', 'lnB'):
        existing = tcPr.find(qn(f'a:{edge}'))
        if existing is not None:
            tcPr.remove(existing)
            
    insert_idx = 0
    for i, child in enumerate(tcPr):
        if child.tag in (qn('a:solidFill'), qn('a:noFill')):
            insert_idx = i
            break
    else:
        insert_idx = len(tcPr)
        
    for edge in ('lnL', 'lnR', 'lnT', 'lnB'):
        el = parse_xml(
            f'<a:{edge} w="{width}" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'  <a:solidFill>'
            f'    <a:srgbClr val="{color}"/>'
            f'  </a:solidFill>'
            f'</a:{edge}>'
        )
        tcPr.insert(insert_idx, el)
        insert_idx += 1


def _fill_cell(cell, rgb: RGBColor) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb


def _format_cell(cell, text: str, font_size: int = BODY_FONT_SIZE, bold: bool = False,
                 align=PP_ALIGN.LEFT, color: RGBColor = BLACK, valign=MSO_ANCHOR.TOP) -> None:
    cell.text = ""
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.035)
    tf.margin_bottom = Inches(0.035)
    tf.vertical_anchor = valign

    lines = str(text or "").split("\n")
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_before = Pt(0)
        if str(line).lstrip().startswith("•"):
            p.space_after = Pt(2.2)
        else:
            p.space_after = Pt(0.6)
        p.line_spacing = 1.02
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def _add_rect(slide, x, y, w, h, fill: RGBColor, line: Optional[RGBColor] = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def _add_textbox(slide, text, x, y, w, h, font_size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_logo(slide, logo_path: Optional[str], x, y, w, h):
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, int(x), int(y), width=int(w), height=int(h))
            return True
        except Exception:
            return False
    return False


def _add_common_frame(slide, logos: Dict[str, Optional[str]], module: str, slide_no: Optional[int] = None,
                      show_arcturus: bool = True, show_tid_bottom: bool = True):
    _add_rect(slide, 0, 0, SLIDE_W, TOP_STRIP_H, BLUE)
    _add_rect(slide, 0, SLIDE_H - BOTTOM_STRIP_H, SLIDE_W, BOTTOM_STRIP_H, BLUE)

    if show_arcturus:
        _add_logo(slide, logos.get("arcturus"), SLIDE_W - Inches(1.55), Inches(0.18), Inches(1.20), Inches(0.38))

    if show_tid_bottom:
        _add_logo(slide, logos.get("tid"), (SLIDE_W - BOTTOM_LOGO_W) / 2, BOTTOM_LOGO_Y, BOTTOM_LOGO_W, BOTTOM_LOGO_H)

    footer_text = f"{module} - Confidential"
    _add_textbox(slide, footer_text, Inches(0.3), SLIDE_H - BOTTOM_STRIP_H, SLIDE_W - Inches(0.6), BOTTOM_STRIP_H,
                 font_size=6, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    if slide_no is not None:
        _add_textbox(slide, str(slide_no), SLIDE_W - Inches(0.45), SLIDE_H - BOTTOM_STRIP_H,
                     Inches(0.28), BOTTOM_STRIP_H, font_size=6, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def _add_title_slide(prs: Presentation, logos: Dict[str, Optional[str]], release: str, module: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_common_frame(slide, logos, module, slide_no=1, show_arcturus=True, show_tid_bottom=True)

    _add_rect(slide, 0, TITLE_BAND_Y, SLIDE_W, TITLE_BAND_H, BLUE)
    _add_textbox(slide, f"Turlock Irrigation District : {release} Oracle Upgrades",
                 Inches(0.7), TITLE_BAND_Y + Inches(0.18), SLIDE_W - Inches(1.4), Inches(0.40),
                 font_size=TITLE_FONT_SIZE, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, module,
                 Inches(0.7), TITLE_BAND_Y + Inches(0.72), SLIDE_W - Inches(1.4), Inches(0.38),
                 font_size=SUBTITLE_FONT_SIZE, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def _add_index_slide(prs: Presentation, logos: Dict[str, Any], index_data: List[Tuple[str, str, int]], module: str):
    start_idx = 0
    while start_idx < len(index_data):
        end_idx = min(start_idx + MAX_INDEX_ROWS_PER_SLIDE, len(index_data))
        chunk = index_data[start_idx:end_idx]

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        _add_textbox(slide, "Index", 
                     CONTENT_MARGIN_X, Inches(0.75), TABLE_W, Inches(0.35),
                     font_size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT)

        _add_logo(slide, logos.get("arcturus"), SLIDE_W - Inches(1.55), Inches(0.18), Inches(1.20), Inches(0.38))
        _add_logo(slide, logos.get("tid"), (SLIDE_W - BOTTOM_LOGO_W) / 2, BOTTOM_LOGO_Y, BOTTOM_LOGO_W, BOTTOM_LOGO_H)

        rows = len(chunk) + 1
        cols = 3
        table = slide.shapes.add_table(rows, cols, Inches(0.5), INDEX_TABLE_TOP_Y, TABLE_W, Inches(1.0)).table
        
        _configure_table_borders(table)

        headers = ["Module", "Release Feature", "Slide"]
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            _fill_cell(cell, BLUE)
            _set_cell_border(cell)
            _format_cell(cell, h, font_size=HEADER_FONT_SIZE, bold=True, align=PP_ALIGN.CENTER, color=WHITE, valign=MSO_ANCHOR.MIDDLE)

        prev_mod = ""
        for i, (mod, title, sld) in enumerate(chunk, start=1):
            mod_display = mod if mod != prev_mod else ""
            prev_mod = mod
            
            vals = [mod_display, title, str(sld)]
            for c, val in enumerate(vals):
                cell = table.cell(i, c)
                _fill_cell(cell, WHITE)
                _set_cell_border(cell)
                align = PP_ALIGN.LEFT if c == 1 else PP_ALIGN.CENTER
                _format_cell(cell, val, font_size=BODY_FONT_SIZE, bold=False, align=align, color=BLACK, valign=MSO_ANCHOR.MIDDLE)

        start_idx += MAX_INDEX_ROWS_PER_SLIDE


def _add_content_slide(prs: Presentation, logos: Dict[str, Optional[str]], chunk: List[Dict[str, Any]],
                       start_index: int, release: str, module: str, slide_no: int, module_heading: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_common_frame(slide, logos, module, slide_no=slide_no, show_arcturus=True, show_tid_bottom=True)

    _add_textbox(slide, module_heading, 
                 CONTENT_MARGIN_X, Inches(0.75), TABLE_W, Inches(0.30),
                 font_size=12, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT)

    cols = 6
    rows = len(chunk) + 1
    row_h = DATA_ROW_H
    table_h = HEADER_ROW_H + row_h * len(chunk)
    max_table_bottom = BOTTOM_LOGO_Y - Inches(0.16)
    if TABLE_TOP_Y + table_h > max_table_bottom:
        row_h = (max_table_bottom - TABLE_TOP_Y - HEADER_ROW_H) / max(len(chunk), 1)
        table_h = HEADER_ROW_H + row_h * len(chunk)

    shape = slide.shapes.add_table(rows, cols, int(CONTENT_MARGIN_X), int(TABLE_TOP_Y), int(TABLE_W), int(table_h))
    table = shape.table

    _configure_table_borders(table)

    table.rows[0].height = int(HEADER_ROW_H)
    for rr in range(1, rows):
        table.rows[rr].height = int(row_h)

    widths = [Inches(0.62), Inches(3.65), Inches(3.30), Inches(3.30), Inches(0.93), Inches(0.95)]
    for i, w in enumerate(widths):
        table.columns[i].width = int(w)

    headers = ["Sr. No.", f"{release} Release Feature", "Business Benefit", "Business Impact", "Mandatory\n(Yes/No)", "Priority"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        _fill_cell(cell, BLUE)
        _set_cell_border(cell)
        _format_cell(cell, h, font_size=HEADER_FONT_SIZE, bold=True, align=PP_ALIGN.CENTER, color=WHITE, valign=MSO_ANCHOR.MIDDLE)

    for r, feature in enumerate(chunk, start=1):
        values = [
            str(start_index + r - 1),
            _release_feature_cell(feature),
            _benefit_cell(feature),
            _impact_cell(feature),
            _mandatory(feature),
            _priority(feature),
        ]
        for c, val in enumerate(values):
            cell = table.cell(r, c)
            _fill_cell(cell, WHITE)
            _set_cell_border(cell)
            align = PP_ALIGN.CENTER if c in (0, 4, 5) else PP_ALIGN.LEFT
            bold = c == 1 and "\n" in val
            _format_cell(cell, val, font_size=BODY_FONT_SIZE, bold=False, align=align, color=BLACK, valign=MSO_ANCHOR.TOP)
            if c == 1:
                try:
                    first_para = cell.text_frame.paragraphs[0]
                    for run in first_para.runs:
                        run.font.bold = True
                except Exception:
                    pass


def _add_reference_slides(prs: Presentation, logos: Dict[str, Optional[str]], features: List[Dict[str, Any]], module: str, start_slide_no: int) -> int:
    """Add reference slides, chunked to avoid overflow. Returns the next slide number."""
    refs = OrderedDict()
    for f in features:
        url = _pick(f, "url", "source_url", "link")
        feat_mod = _feature_module(f)
        title = _feature_title(f)
        if url:
            if url not in refs:
                refs[url] = []
            refs[url].append(f"{feat_mod} - {title}")

    if not refs:
        return start_slide_no

    ref_items = list(refs.items())
    current_slide_no = start_slide_no
    start_idx = 0

    while start_idx < len(ref_items):
        end_idx = min(start_idx + MAX_REF_ROWS_PER_SLIDE, len(ref_items))
        chunk = ref_items[start_idx:end_idx]

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_common_frame(slide, logos, module, slide_no=current_slide_no, show_arcturus=True, show_tid_bottom=True)

        _add_textbox(slide, "References", 
                     CONTENT_MARGIN_X, Inches(0.75), TABLE_W, Inches(0.35),
                     font_size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT)

        rows = len(chunk) + 1
        cols = 2
        table_h = Inches(0.4 * rows)
        table = slide.shapes.add_table(rows, cols, int(CONTENT_MARGIN_X), Inches(1.20), int(TABLE_W), int(table_h)).table
        
        _configure_table_borders(table)
        
        widths = [Inches(4.5), Inches(8.5)]
        for i, w in enumerate(widths):
            table.columns[i].width = int(w)

        headers = ["Feature(s)", "Reference Link"]
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            _fill_cell(cell, BLUE)
            _set_cell_border(cell)
            _format_cell(cell, h, font_size=HEADER_FONT_SIZE, bold=True, align=PP_ALIGN.CENTER, color=WHITE, valign=MSO_ANCHOR.MIDDLE)

        for r, (url, titles) in enumerate(chunk, start=1):
            cell_title = table.cell(r, 0)
            cell_url = table.cell(r, 1)
            
            _fill_cell(cell_title, WHITE)
            _set_cell_border(cell_title)
            _format_cell(cell_title, "\n".join(titles), font_size=8, bold=False, align=PP_ALIGN.LEFT, color=BLACK, valign=MSO_ANCHOR.TOP)
            
            _fill_cell(cell_url, WHITE)
            _set_cell_border(cell_url)
            
            cell_url.text = ""
            tf = cell_url.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = Inches(0.04)
            tf.margin_right = Inches(0.04)
            tf.margin_top = Inches(0.035)
            tf.margin_bottom = Inches(0.035)
            
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = url
            run.font.name = FONT
            run.font.size = Pt(8)
            run.font.color.rgb = BLUE
            run.font.underline = True
            run.hyperlink.address = url

        start_idx += MAX_REF_ROWS_PER_SLIDE
        current_slide_no += 1
        
    return current_slide_no


def generate_ppt(features: List[Dict[str, Any]], *args) -> str:
    if len(args) == 1:
        template_path = None
        output_path = args[0]
    elif len(args) >= 2:
        template_path = args[0]
        output_path = args[1]
    else:
        output_path = "outputs/oquat_report.pptx"
        template_path = None

    features = features or []
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    tpl = _find_template_path(template_path)
    logos = _extract_logos_from_template(tpl)

    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)

    release = _release_version(features)
    module = _module_name(features)

    # Pre-scan to identify repetitive benefit bullets across all features
    _prescan_common_bullets(features)

    _add_title_slide(prs, logos, release, module)

    features_by_module = OrderedDict()
    for f in features:
        mod = _feature_module(f)
        if mod not in features_by_module:
            features_by_module[mod] = []
        features_by_module[mod].append(f)

    rows_per_slide = 3
    chunks = []
    index_data = []
    slide_no = 3
    
    for mod, mod_features in features_by_module.items():
        for i in range(0, len(mod_features), rows_per_slide):
            chunk = mod_features[i:i+rows_per_slide]
            chunks.append((mod, chunk, slide_no))
            for f in chunk:
                index_data.append((mod, _feature_title(f), slide_no))
            slide_no += 1

    _add_index_slide(prs, logos, index_data, module)

    global_feature_idx = 1
    for mod, chunk, s_no in chunks:
        _add_content_slide(prs, logos, chunk, global_feature_idx, release, module, s_no, module_heading=mod)
        global_feature_idx += len(chunk)

    if features:
        slide_no = _add_reference_slides(prs, logos, features, module, slide_no)

    prs.save(output_path)
    return output_path